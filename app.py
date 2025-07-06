from dotenv import load_dotenv
import os
import json
import shutil
from datetime import datetime
from flask import Flask, render_template, request, send_file, redirect, url_for, flash, jsonify
from werkzeug.utils import secure_filename
import docx
import fitz  # PyMuPDF
from pptx import Presentation
from config import llm_config, LLMProvider
from llm_service import llm_service
import logging
from logging.handlers import RotatingFileHandler
import traceback

# --- CARICA CONFIG ---
load_dotenv()

# --- CONFIG FLASK ---
app = Flask(__name__)
app.config['SECRET_KEY'] = 'a-super-secret-key-for-flash-messages'
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ARCHIVE_FOLDER'] = 'archive'
app.config['ALLOWED_EXTENSIONS'] = {'txt', 'pdf', 'docx'}
app.config['TEMPLATE_PATH'] = os.path.join('static', 'template.pptx')

# Logging su file
if not os.path.exists('logs'):
    os.makedirs('logs')
file_handler = RotatingFileHandler('logs/app.log', maxBytes=10240, backupCount=3)
file_handler.setLevel(logging.ERROR)
file_handler.setFormatter(logging.Formatter('%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'))
app.logger.addHandler(file_handler)

# Handler globale per errori 500
@app.errorhandler(500)
def internal_error(error):
    tb = traceback.format_exc()
    app.logger.error(f"Errore 500: {error}\nTraceback:\n{tb}")
    return render_template('500.html', error=error, traceback=tb), 500

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['ARCHIVE_FOLDER'], exist_ok=True)

# --- UTILITY ---
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def create_session_folder(first_filename):
    """Crea una cartella di sessione con naming convention"""
    # Rimuove estensione e caratteri speciali dal nome del primo file
    base_name = os.path.splitext(first_filename)[0]
    safe_name = "".join(c for c in base_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
    safe_name = safe_name.replace(' ', '_')[:50]  # Limita lunghezza
    
    # Aggiunge timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    folder_name = f"{safe_name}_{timestamp}"
    
    # Crea il percorso completo
    session_path = os.path.join(app.config['ARCHIVE_FOLDER'], folder_name)
    os.makedirs(session_path, exist_ok=True)
    
    return session_path, folder_name

def save_files_to_session(files, session_path):
    """Salva i file di input nella cartella di sessione"""
    saved_files = []
    for file in files:
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(session_path, filename)
            file.save(filepath)
            saved_files.append(filepath)
    return saved_files

def create_session_presentation(slides_content, session_path, session_name):
    """Crea la presentazione nella cartella di sessione"""
    prs = Presentation(app.config['TEMPLATE_PATH'])
    for slide_text in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_text.get("title", "")
        slide.placeholders[1].text = slide_text.get("content", "")
    
    output_filename = f"{session_name}_presentation.pptx"
    output_path = os.path.join(session_path, output_filename)
    prs.save(output_path)
    return output_path

def extract_text_from_file(filepath):
    ext = filepath.rsplit('.', 1)[1].lower()
    text = ""
    if ext == 'txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
    elif ext == 'pdf':
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text()
    elif ext == 'docx':
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            text += para.text + '\n'
    return text

def create_presentation(slides_content):
    prs = Presentation(app.config['TEMPLATE_PATH'])
    for slide_text in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_text.get("title", "")
        slide.placeholders[1].text = slide_text.get("content", "")
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], "output.pptx")
    prs.save(output_path)
    return output_path

def generate_slide_content(prompt_text):
    try:
        prompt = f"Analizza il seguente contenuto e genera una struttura di slide PowerPoint professionale:\n\n{prompt_text}\n\nRispondi SOLO con un JSON valido contenente una lista di oggetti con 'title' e 'content' per ogni slide."
        response = llm_service.generate_content(prompt)
        
        # Prova a pulire e parsificare il JSON
        try:
            # Rimuove eventuali backticks e testo extra
            json_start = response.find('[')
            json_end = response.rfind(']') + 1
            if json_start != -1 and json_end > json_start:
                json_str = response[json_start:json_end]
                slides_content = json.loads(json_str)
            else:
                # Se non trova JSON, prova a parsificare direttamente
                slides_content = json.loads(response)
        except json.JSONDecodeError:
            # Se il parsing fallisce, crea slide con contenuto grezzo
            slides_content = [
                {"title": "Contenuto Generato", "content": response},
                {"title": "Nota", "content": "Il modello non ha restituito un JSON valido. Contenuto mostrato in forma grezza."}
            ]
        
        # Verifica che sia una lista di dizionari
        if not isinstance(slides_content, list):
            slides_content = [{"title": "Contenuto", "content": str(slides_content)}]
        
        # Verifica che ogni elemento abbia title e content
        for slide in slides_content:
            if not isinstance(slide, dict):
                slide = {"title": "Slide", "content": str(slide)}
            if 'title' not in slide:
                slide['title'] = "Slide"
            if 'content' not in slide:
                slide['content'] = ""
        
        return slides_content
    except Exception as e:
        return [{"title": "Errore generazione", "content": f"Errore: {str(e)}"}]

# --- ROUTES ---
@app.route("/", methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Validazione base
        if 'file' not in request.files:
            flash('Nessun file selezionato')
            return redirect(request.url)
        
        files = request.files.getlist('file')
        if not files or all(f.filename == '' for f in files):
            flash('Fornire almeno 1 file di input in uno dei formati validi (PDF, DOCX, TXT)')
            return redirect(request.url)
        
        # Filtra solo i file validi
        valid_files = [f for f in files if f and f.filename != '' and allowed_file(f.filename)]
        
        if not valid_files:
            flash('Fornire almeno 1 file di input in uno dei formati validi (PDF, DOCX, TXT)')
            return redirect(request.url)
        
        try:
            # Crea cartella di sessione basata sul primo file
            first_filename = secure_filename(valid_files[0].filename)
            session_path, session_name = create_session_folder(first_filename)
            
            # Salva tutti i file nella cartella di sessione
            saved_files = save_files_to_session(valid_files, session_path)
            
            # Estrai testo da tutti i file
            texts = []
            for filepath in saved_files:
                text = extract_text_from_file(filepath)
                if text.strip():  # Solo se il file ha contenuto
                    texts.append(text)
            
            if not texts:
                flash('I file caricati non contengono testo leggibile')
                # Rimuovi cartella vuota
                shutil.rmtree(session_path, ignore_errors=True)
                return redirect(request.url)
            
            # Genera le slide
            combined_text = '\n\n--- NUOVO DOCUMENTO ---\n\n'.join(texts)
            slides_content = generate_slide_content(combined_text)
            
            # Crea presentazione nella cartella di sessione
            pptx_path = create_session_presentation(slides_content, session_path, session_name)
            
            # Crea anche una copia temporanea per il download immediato
            temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{session_name}_presentation.pptx")
            shutil.copy2(pptx_path, temp_path)
            
            return send_file(temp_path, as_attachment=True, download_name=f"{session_name}_presentation.pptx")
            
        except Exception as e:
            app.logger.error(f"Errore nella generazione: {str(e)}")
            flash(f'Errore nella generazione della presentazione: {str(e)}')
            return redirect(request.url)
    
    return render_template('index.html')

@app.route("/config")
def config_page():
    current_model_config = llm_config.get_current_model()
    current_model_name = current_model_config.name if current_model_config else "Nessuno"
    
    # Modelli cloud disponibili solo se c'è la key
    available_cloud = llm_service.available_cloud_models()
    # Modelli locali effettivamente disponibili
    available_local = llm_service.available_local_models()
    # Per ogni provider cloud, indica se manca la key
    missing_keys = {
        'openai': not bool(llm_config.get_api_key(LLMProvider.OPENAI)),
        'anthropic': not bool(llm_config.get_api_key(LLMProvider.ANTHROPIC)),
        'google': not bool(llm_config.get_api_key(LLMProvider.GOOGLE)),
    }
    return render_template('config.html',
        current_model=current_model_name,
        current_model_id=llm_config.current_model,
        available_cloud=available_cloud,
        available_local=available_local,
        missing_keys=missing_keys,
        local_backend=llm_config.local_backend,
        local_endpoint=llm_config.local_endpoint,
        llm_config=llm_config)

@app.route("/api/set_model", methods=['POST'])
def set_model():
    data = request.get_json()
    model_id = data.get('model_id')
    if model_id and model_id in llm_config.models:
        llm_config.set_current_model(model_id)
        return jsonify({"status": "success", "message": f"Modello impostato su {llm_config.models[model_id].name}"})
    return jsonify({"status": "error", "message": "Modello non trovato"})

@app.route("/api/update_model_params", methods=['POST'])
def update_model_params():
    data = request.get_json()
    model_id = data.get('model_id')
    params = data.get('params', {})
    
    if model_id and model_id in llm_config.models:
        llm_config.update_model_params(model_id, **params)
        return jsonify({"status": "success", "message": "Parametri aggiornati"})
    return jsonify({"status": "error", "message": "Modello non trovato"})

@app.route("/api/set_api_key", methods=['POST'])
def set_api_key():
    data = request.get_json()
    provider_name = data.get('provider')
    api_key = data.get('api_key')
    
    try:
        provider = LLMProvider(provider_name)
        llm_config.set_api_key(provider, api_key)
        # Salva immediatamente la configurazione per persistenza
        llm_config.save_config()
        # Reinizializza il servizio LLM
        llm_service._initialize_clients()
        return jsonify({"status": "success", "message": f"Chiave API per {provider_name} impostata e salvata"})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

@app.route("/api/test_connection", methods=['POST'])
def test_connection():
    data = request.get_json()
    provider_name = data.get('provider')
    
    try:
        provider = LLMProvider(provider_name)
        result = llm_service.test_connection(provider)
        return jsonify(result)
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

@app.route("/api/get_available_models", methods=['POST'])
def get_available_models():
    data = request.get_json()
    provider_name = data.get('provider')
    
    try:
        provider = LLMProvider(provider_name)
        result = llm_service.get_available_models(provider)
        return jsonify(result)
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

@app.route("/api/set_local_backend", methods=['POST'])
def set_local_backend():
    data = request.get_json()
    backend = data.get('backend')
    endpoint = data.get('endpoint')
    if backend in ['ollama', 'lmstudio'] and endpoint:
        llm_config.set_local_backend(backend, endpoint)
        return jsonify({"status": "success", "message": "Backend locale aggiornato"})
    return jsonify({"status": "error", "message": "Dati non validi"})

@app.route("/api/get_config", methods=['GET'])
def get_config():
    # Restituisco solo i modelli effettivamente disponibili
    available_cloud = llm_service.available_cloud_models()
    available_local = llm_service.available_local_models()
    config = {
        'current_model': llm_config.current_model,
        'local_backend': llm_config.local_backend,
        'local_endpoint': llm_config.local_endpoint,
        'available_cloud': {k: v.__dict__ for k, v in available_cloud.items()},
        'available_local': {k: v.__dict__ for k, v in available_local.items()},
        'api_keys': {
            'openai': bool(llm_config.get_api_key(LLMProvider.OPENAI)),
            'anthropic': bool(llm_config.get_api_key(LLMProvider.ANTHROPIC)),
            'google': bool(llm_config.get_api_key(LLMProvider.GOOGLE)),
        }
    }
    return jsonify(config)

@app.route("/api/refresh_models", methods=['GET'])
def refresh_models():
    # Endpoint per refresh dinamico della lista modelli
    available_cloud = llm_service.available_cloud_models()
    available_local = llm_service.available_local_models()
    return jsonify({
        'available_cloud': {k: v.__dict__ for k, v in available_cloud.items()},
        'available_local': {k: v.__dict__ for k, v in available_local.items()}
    })

@app.route("/api/list_models", methods=['POST'])
def api_list_models():
    data = request.get_json()
    provider = data.get('provider')
    result = []
    try:
        if provider == 'openai':
            models = llm_service.get_available_models(LLMProvider.OPENAI)
            if models.get('status') == 'success':
                result = models.get('models', [])
        elif provider == 'anthropic':
            models = llm_service.get_available_models(LLMProvider.ANTHROPIC)
            if models.get('status') == 'success':
                result = models.get('models', [])
        elif provider == 'google':
            models = llm_service.get_available_models(LLMProvider.GOOGLE)
            if models.get('status') == 'success':
                result = models.get('models', [])
        elif provider == 'local':
            models = llm_service.get_available_models(LLMProvider.LOCAL)
            if models.get('status') == 'success':
                result = models.get('models', [])
        return jsonify({'status': 'success', 'models': result})
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)})

@app.route("/api/default_endpoint", methods=['POST'])
def api_default_endpoint():
    data = request.get_json()
    backend = data.get('backend')
    # Endpoint di default per Ollama e LM Studio
    defaults = {
        'ollama': 'http://localhost:11434',
        'lmstudio': 'http://localhost:1234',
    }
    # Recupera l'ultimo endpoint usato per il backend, se presente
    last = getattr(llm_config, f'last_endpoint_{backend}', None)
    endpoint = last or defaults.get(backend, '')
    return jsonify({'status': 'success', 'endpoint': endpoint})

@app.route("/api/get_system_prompt", methods=['GET'])
def get_system_prompt():
    return jsonify({"status": "success", "prompt": llm_config.system_prompt})

@app.route("/api/set_system_prompt", methods=['POST'])
def set_system_prompt():
    data = request.get_json()
    prompt = data.get('prompt', '').strip()
    if not prompt:
        return jsonify({"status": "error", "message": "Prompt non valido"})
    llm_config.set_system_prompt(prompt)
    return jsonify({"status": "success", "message": "Prompt aggiornato"})

# --- START SERVER ---
if __name__ == "__main__":
    app.run(debug=False, port=8080)
